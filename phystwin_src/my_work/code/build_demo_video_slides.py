#!/usr/bin/env python3
"""build_demo_video_slides.py — three standalone slides, each side-by-side
video comparison.

Slide A : descriptor effect, isolated     (rope replay, single_lift_rope)
          BC base  (err 0.36)  vs  BC + 2-d desc  (err 0.13)

Slide B : RL effect on top of BC + desc   (rope ramp, single_push_rope_1)
          BC + desc  (fall 0.30)  vs  RL + desc  (fall 0.58)

Slide C : Full pipeline impact            (rope ramp, single_push_rope_1)
          BC base  (fall 0.04, doesn't retract)  vs  RL + desc  (fall 0.58)

Each video's poster image is its OWN first frame, so the static thumb
shows what the video actually looks like (no out-of-context cartoon).
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
VID  = ROOT / "presentation_results" / "videos"
FIG  = ROOT / "presentation_results" / "figures"
OUT  = ROOT / "presentation_results" / "demo_video_slides.pptx"

INK    = RGBColor(0x11, 0x18, 0x27)
DIM    = RGBColor(0x4b, 0x55, 0x63)
ACCENT = RGBColor(0x16, 0xa3, 0x4a)
BLUE   = RGBColor(0x25, 0x63, 0xeb)
GRAY   = RGBColor(0x6b, 0x72, 0x80)
BG     = RGBColor(0xff, 0xff, 0xff)


def add_text(slide, l, t, w, h, txt, *,
              size=18, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
        r.font.name = "Calibri"


def title_block(slide, title, subtitle):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
              title, size=30, bold=True)
    add_text(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
              subtitle, size=15, color=DIM, italic=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.45), Inches(12.3), Emu(12700))
    rule.fill.solid(); rule.fill.fore_color.rgb = DIM
    rule.line.fill.background()


def set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = ""
    for i, line in enumerate(text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for r in p.runs:
            r.font.size = Pt(13)


def two_video_slide(prs, title, subtitle,
                     left_video, left_poster, left_color, left_top, left_bot,
                     right_video, right_poster, right_color, right_top, right_bot,
                     footer_caption, notes):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    title_block(s, title, subtitle)

    # Each video gets its OWN poster = the actual force-tracking graph.
    s.shapes.add_movie(str(left_video),
                        Inches(0.5), Inches(2.0),
                        Inches(6.0), Inches(3.8),
                        poster_frame_image=str(left_poster),
                        mime_type="video/mp4")
    s.shapes.add_movie(str(right_video),
                        Inches(6.8), Inches(2.0),
                        Inches(6.0), Inches(3.8),
                        poster_frame_image=str(right_poster),
                        mime_type="video/mp4")
    add_text(s, Inches(0.5), Inches(5.95), Inches(6.0), Inches(0.4),
              left_top, size=15, bold=True, color=left_color,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.35), Inches(6.0), Inches(0.4),
              left_bot, size=12, color=DIM, italic=True,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.8), Inches(5.95), Inches(6.0), Inches(0.4),
              right_top, size=15, bold=True, color=right_color,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.8), Inches(6.35), Inches(6.0), Inches(0.4),
              right_bot, size=12, color=DIM, italic=True,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
              footer_caption, size=12, italic=True, color=DIM,
              align=PP_ALIGN.CENTER)
    set_notes(s, notes)


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    POSTERS = FIG / "posters"
    # ---------- Slide A : descriptor effect on rope REPLAY ----------
    two_video_slide(
        prs,
        title    = "Contribution 1 · Material descriptor",
        subtitle = "same network, two extra static numbers in the input — rope replay (single_lift_rope)",
        left_video  = VID / "09_BC_rope_replay.mp4",
        left_poster = POSTERS / "09_BC_rope_replay_frame0.png",
        left_color  = GRAY,
        left_top    = "BC  (no descriptor)",
        left_bot    = "err 0.36  ·  tracks shape, but lags",
        right_video = VID / "10_BCdesc_rope_replay.mp4",
        right_poster = POSTERS / "10_BCdesc_rope_replay_frame0.png",
        right_color = BLUE,
        right_top   = "BC + 2-d descriptor",
        right_bot   = "err 0.13  ·  tracks closely  (−62%)",
        footer_caption = "Same MLP. Same supervised loss. The only difference is the [log_Y, log_F] descriptor in the input.",
        notes = """
            ~25 seconds. Both rollouts are on the SAME case
            (single_lift_rope, replay profile). Same supervised
            training pipeline. The only variable is whether the policy
            sees the 2-d material descriptor.

            LEFT (BC base, no descriptor): error ratio 0.36. The
            policy tracks the general shape of the goal force, but
            lags noticeably and undershoots peaks.

            RIGHT (BC + 2-d descriptor): error ratio 0.13 — a 62%
            improvement. Same MLP architecture, same loss; the only
            extra input is two static numbers (log of mean stiffness,
            log of mean recorded force magnitude). With that hint the
            policy tracks the goal closely.

            Both posters show the actual force-tracking curves —
            achieved vs goal magnitude over time. The right poster
            shows tighter alignment.

            One-liner: "Telling the network 'this is rope, force
            scale ~15 kN' — two static numbers in the input —
            cuts tracking error by more than half."
        """,
    )

    # ---------- Slide B : RL effect ON TOP OF descriptor ----------
    two_video_slide(
        prs,
        title    = "Contribution 2 · PPO with BC warm-start",
        subtitle = "isolated RL effect — both already have the descriptor (rope ramp, single_push_rope_1)",
        left_video  = VID / "08_BCdesc_rope_ramp.mp4",
        left_poster = POSTERS / "08_BCdesc_rope_ramp_frame0.png",
        left_color  = BLUE,
        left_top    = "BC + 2-d descriptor",
        left_bot    = "fall fraction 0.30  ·  partial retraction",
        right_video = VID / "01_RL_rope_ramp_HEADLINE.mp4",
        right_poster = POSTERS / "01_RL_rope_ramp_HEADLINE_frame0.png",
        right_color = ACCENT,
        right_top   = "RL + 2-d descriptor",
        right_bot   = "fall fraction 0.58  ·  near-double retraction",
        footer_caption = "Modest but measurable: RL nearly doubles the retraction the descriptor alone provides.",
        notes = """
            ~20 seconds. Both policies already have the material
            descriptor, so the only variable on THIS slide is BC vs
            PPO refinement.

            LEFT (BC + descriptor on rope ramp): rises to ~7800 N
            peak, then partially retracts. Fall fraction 0.30 —
            roughly a third of the peak force decays by the end.

            RIGHT (RL + descriptor on the same case): same rise, but
            retracts further. Fall fraction 0.58 — nearly double the
            BC retraction.

            Be honest here: this is a modest improvement, not a
            transformation. The big "BC barely retracts at all"
            failure is the bare baseline without ramp data or
            descriptor — that's what the next slide shows.

            One-liner: "On top of the descriptor, RL roughly doubles
            the retraction. The bigger story — what RL plus the
            descriptor together accomplish — is on the next slide."

            If asked why the err_ratio change is modest (0.62 → 0.52):
            the rising portion dominates the average and looks the
            same for both. Fall fraction isolates the falling-edge
            behavior — the closed-loop state where they actually
            differ.
        """,
    )

    # ---------- Slide C : Full pipeline impact on rope ramp ----------
    two_video_slide(
        prs,
        title    = "Full pipeline · BC base vs everything stacked",
        subtitle = "what the final system does that the supervised baseline can't — rope ramp",
        left_video  = VID / "11_BC_rope_ramp.mp4",
        left_poster = POSTERS / "11_BC_rope_ramp_frame0.png",
        left_color  = GRAY,
        left_top    = "BC base  (no descriptor, no ramp data)",
        left_bot    = "fall fraction 0.04  ·  doesn’t retract at all  ·  final 8.7 kN",
        right_video = VID / "01_RL_rope_ramp_HEADLINE.mp4",
        right_poster = POSTERS / "01_RL_rope_ramp_HEADLINE_frame0.png",
        right_color = ACCENT,
        right_top   = "RL + 2-d descriptor",
        right_bot   = "fall fraction 0.58  ·  retracts to 3.3 kN  ·  goal end 0 kN",
        footer_caption = "Stacked contributions: synthetic ramp data + material descriptor + RL refinement.",
        notes = """
            ~20 seconds. This slide is the FULL system delta — bare
            baseline on the left, complete pipeline on the right.

            LEFT (BC base, Step 2): no material descriptor, trained
            without the synthetic ramp augmentation. When the goal
            force drops back to zero at the end, the policy keeps
            pushing — the final force is 8.7 kN, almost the same as
            the peak. Fall fraction 0.04.

            RIGHT (RL + descriptor, full pipeline): retracts to 3.3 kN
            by the end. Fall fraction 0.58.

            Honest framing: this is descriptor + ramp data + RL
            stacked, NOT a clean RL-only delta. The previous slide
            shows the isolated RL contribution (0.30 → 0.58). Most of
            the retraction was unlocked by the descriptor + ramp data;
            RL refined it further.

            One-liner: "The end-to-end win is real — bare BC doesn't
            retract, the full system does. The credit is shared
            across three improvements; RL is the smallest of them but
            still adds measurable lift."

            If anyone asks "which contribution mattered most" — the
            material descriptor; details on the bar chart slide.
        """,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
