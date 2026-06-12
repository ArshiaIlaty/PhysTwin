#!/usr/bin/env python3
"""build_rl_training_slide.py — single slide explaining RL training.

Writes ONE slide into a standalone .pptx so it can be dragged into the
main deck without regenerating everything. Same light theme palette as
build_presentation_pptx.py.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT_PPTX = ROOT / "presentation_results" / "rl_training_slide.pptx"

INK    = RGBColor(0x11, 0x18, 0x27)
DIM    = RGBColor(0x4b, 0x55, 0x63)
ACCENT = RGBColor(0x16, 0xa3, 0x4a)
BLUE   = RGBColor(0x25, 0x63, 0xeb)
ORANGE = RGBColor(0xc2, 0x41, 0x0c)
BG     = RGBColor(0xff, 0xff, 0xff)


def add_text(slide, left, top, width, height, text,
              size=20, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_round_box(slide, x, y, w, h, fill, label_top, label_bot=None):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.color.rgb = INK
    r.line.width = Pt(1)
    tf = r.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = label_top
    run.font.size = Pt(13); run.font.bold = True; run.font.name = "Calibri"
    run.font.color.rgb = INK
    if label_bot:
        para2 = tf.add_paragraph(); para2.alignment = PP_ALIGN.CENTER
        run2 = para2.add_run(); run2.text = label_bot
        run2.font.size = Pt(10); run2.font.color.rgb = DIM
        run2.font.name = "Calibri"


def add_arrow(slide, x1, y1, x2, y2, color=INK):
    arr = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    arr.line.color.rgb = color; arr.line.width = Pt(2)
    # add tail-to-head arrow marker
    from pptx.oxml.ns import qn
    ln = arr.line._get_or_add_ln()
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        from pptx.oxml.ns import nsmap
        from lxml import etree
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med'); tailEnd.set('len', 'med')


def set_notes(slide, text: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = ""
    for i, line in enumerate(text.strip().split("\n")):
        para = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(13)


def title_block(slide, title, subtitle):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
              title, size=32, bold=True, color=INK)
    add_text(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
              subtitle, size=16, color=DIM, italic=True)
    # thin rule under title
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.45), Inches(12.3), Emu(12700))
    rule.fill.solid(); rule.fill.fore_color.rgb = DIM
    rule.line.fill.background()


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

    title_block(s, "How we trained the RL policy",
                  "small steps from the BC policy, scored by closed-loop force tracking")

    # ---- 4-box training loop along the top ----
    y, h = 2.0, 1.4
    boxes = [
        ("1.  Start from BC",       "warm-start the actor\nwith BC weights",        "#dbeafe"),
        ("2.  Try variations",       "add small random noise\nto each action",       "#fcd34d"),
        ("3.  Score the rollout",    "−‖achieved − goal‖²\nover the trajectory",     "#fecaca"),
        ("4.  Nudge weights",        "toward higher-scoring\nvariants (KL-clipped)", "#a7f3d0"),
    ]
    x = 0.5; w = 2.85; gap = 0.27
    coords = []
    for label_top, label_bot, color in boxes:
        add_round_box(s, Inches(x), Inches(y), Inches(w), Inches(h),
                       RGBColor.from_string(color[1:]), label_top, label_bot)
        coords.append((x, y, w, h))
        x += w + gap
    # arrows between boxes
    for i in range(len(boxes) - 1):
        x1 = coords[i][0] + coords[i][2]
        x2 = coords[i+1][0]
        add_arrow(s, x1, y + h/2, x2, y + h/2)
    # loop-back arrow under all 4 boxes
    loop_y = y + h + 0.45
    add_arrow(s, coords[-1][0] + coords[-1][2]/2, y + h,
                coords[-1][0] + coords[-1][2]/2, loop_y, color=ACCENT)
    add_arrow(s, coords[-1][0] + coords[-1][2]/2, loop_y,
                coords[0][0]  + coords[0][2]/2,  loop_y, color=ACCENT)
    add_arrow(s, coords[0][0]  + coords[0][2]/2, loop_y,
                coords[0][0]  + coords[0][2]/2, y + h, color=ACCENT)
    add_text(s, Inches(4.0), Inches(loop_y + 0.05), Inches(5.0), Inches(0.35),
              "repeat ~25 times  (50,000 simulator steps total)",
              size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # ---- bottom-left: what makes the loop stable ----
    sec_y = 5.0
    add_text(s, Inches(0.6), Inches(sec_y), Inches(6.0), Inches(0.4),
              "What keeps it from breaking the BC policy:",
              size=15, bold=True, color=INK)
    bullets = [
        "Critic warm-up — the helper value network trains first, so it doesn't feed garbage scores to the actor.",
        "Tiny actor learning rate (3·10⁻⁵) — don't undo the BC weights in one step.",
        "KL clip — cap how much each update can change the policy.",
    ]
    for i, b in enumerate(bullets):
        add_text(s, Inches(0.85), Inches(sec_y + 0.45 + i * 0.45),
                  Inches(6.0), Inches(0.4), "• " + b,
                  size=12, color=DIM)

    # ---- bottom-right: what we trained on ----
    add_text(s, Inches(7.2), Inches(sec_y), Inches(5.6), Inches(0.4),
              "What PPO actually rolled out on:",
              size=15, bold=True, color=INK)
    add_text(s, Inches(7.45), Inches(sec_y + 0.45),
              Inches(5.5), Inches(1.4),
              "•  single_push_rope_1   ·  ramp\n"
              "•  double_lift_cloth_3  ·  ramp\n"
              "•  double_stretch_sloth ·  ramp",
              size=13, color=INK)
    add_text(s, Inches(7.45), Inches(sec_y + 1.55),
              Inches(5.5), Inches(0.4),
              "3 cases  ·  one per material  ·  ramp profile only",
              size=12, italic=True, color=DIM)

    set_notes(s, """
        Walk through the four boxes left-to-right, then talk about the
        loop-back arrow.

        Box 1: We start the RL policy from the supervised BC weights —
        not from random. So even before PPO touches anything, the
        actor is already a competent inverse-dynamics policy.

        Box 2: To learn anything, the policy has to try variations of
        its current behavior. We add small Gaussian noise to each
        action, so each rollout is slightly different from the BC one.

        Box 3: For each rollout, we get one score — how well did the
        achieved force track the goal force in closed loop. Lower
        error = higher score.

        Box 4: We compute a gradient that nudges the weights toward
        higher-scoring variants and away from worse ones. The "KL
        clip" caps how much the policy can change in one update, so
        we don't break the BC initialization in a single bad step.

        Then we loop: 50,000 simulator steps total, ~25 PPO updates,
        about 3 hours on one GPU.

        The bottom-left bullets are the three safety tricks that make
        BC warm-start actually work. Without critic warm-up, the
        random value network in step 3 produces garbage scores in the
        first updates, and the actor gets destroyed.

        The bottom-right list is important honesty: PPO only ever
        rolled out on three cases, one per material, ramp profile.
        The actor's MLP saw every case via the BC warm-start, but the
        PPO gradient signal only came from these three rollouts.

        One-liner to deliver: "RL took the BC policy, added small
        random nudges, scored each rollout in the simulator, and
        bumped the weights toward variants that tracked the force
        better — with three different brakes so the BC behavior
        didn't fall apart in the first few updates."
    """)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX}")


if __name__ == "__main__":
    main()
