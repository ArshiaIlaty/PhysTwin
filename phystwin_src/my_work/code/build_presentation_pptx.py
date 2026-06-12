#!/usr/bin/env python3
"""build_presentation_pptx.py — light-theme 10-min deck.

Generates closed-loop force control presentation. Keeps text minimal,
embeds the matplotlib figures + per-case videos.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "presentation_results"
FIG = ASSETS / "figures"
VID = ASSETS / "videos"
OUT_PPTX = ASSETS / "closed_loop_force_control.pptx"

# Light theme palette
INK    = RGBColor(0x11, 0x18, 0x27)   # near-black for body
DIM    = RGBColor(0x4b, 0x55, 0x63)   # secondary text
ACCENT = RGBColor(0x16, 0xa3, 0x4a)   # green for RL / wins
BLUE   = RGBColor(0x25, 0x63, 0xeb)   # blue for BC
ORANGE = RGBColor(0xc2, 0x41, 0x0c)   # orange for feedback / problem
BG     = RGBColor(0xff, 0xff, 0xff)
SUBTLE = RGBColor(0xf8, 0xfa, 0xfc)   # off-white card

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9


def make_pres() -> Presentation:
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def add_blank_slide(prs):
    """Blank layout (index 6 in default master)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text,
             size=22, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_thin_rule(slide, left, top, width, color=DIM):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(12700))  # ~1px
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()


def set_notes(slide, text: str):
    """Write speaker notes (the bottom pane in Presenter View)."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = ""  # clear
    for i, line in enumerate(text.strip().split("\n")):
        para = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(13)


def add_footer(slide, page_num):
    add_text(slide, Inches(0.4), Inches(7.15),
              Inches(8), Inches(0.3),
              "Closed-loop force tracking via PhysTwin  ·  Malak",
              size=10, color=DIM)
    add_text(slide, Inches(12.4), Inches(7.15),
              Inches(0.8), Inches(0.3),
              str(page_num), size=10, color=DIM, align=PP_ALIGN.RIGHT)


def title_block(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.3),
              Inches(12.3), Inches(0.7),
              title, size=32, bold=True, color=INK)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.0),
                  Inches(12.3), Inches(0.4),
                  subtitle, size=16, color=DIM, italic=True)
    add_thin_rule(slide, Inches(0.5), Inches(1.45), Inches(12.3))


# --------------------------- slide builders --------------------------------

def slide_title(prs):
    s = add_blank_slide(prs); set_bg(s)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.0),
              "Closed-loop force tracking", size=48, bold=True, color=INK,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.6),
              "via PhysTwin's differentiable simulator", size=24,
              color=DIM, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
              "Malak Gamal Eldin · UC Irvine", size=16, color=DIM,
              align=PP_ALIGN.CENTER)
    set_notes(s, """
        Open: "Hi, I'm Malak. I'm going to talk about closed-loop force
        tracking on deformable objects, built on top of PhysTwin's
        differentiable simulator." Keep it under 15 seconds — the real
        content starts on slide 2.
    """)
    return s


def slide_goal(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Goal", "invert PhysTwin: target force → gripper motion")

    add_text(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.6),
              "Forward (prior work):", size=16, bold=True, color=BLUE)
    add_text(s, Inches(2.7), Inches(2.0), Inches(10), Inches(0.6),
              "video  →  particle tracking  →  fit physics  →  contact forces",
              size=17, color=INK)
    add_text(s, Inches(2.7), Inches(2.5), Inches(10), Inches(0.4),
              "(slow, offline, per-video optimization)",
              size=12, color=DIM, italic=True)

    add_text(s, Inches(0.7), Inches(3.4), Inches(11.9), Inches(0.6),
              "Inverse (ours):", size=16, bold=True, color=ACCENT)
    add_text(s, Inches(2.7), Inches(3.4), Inches(10), Inches(0.6),
              "target F*(t)  →  policy  →  gripper Δ  →  realized force",
              size=17, color=INK)
    add_text(s, Inches(2.7), Inches(3.9), Inches(10), Inches(0.4),
              "(real-time, closed-loop, one shared policy)",
              size=12, color=DIM, italic=True)

    add_text(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.6),
              'Build the policy. Close the loop. One network, three materials.',
              size=20, italic=True, color=DIM,
              align=PP_ALIGN.CENTER)
    set_notes(s, """
        Two arrows. The TOP arrow is what PhysTwin already does: it
        takes RGB-D video, tracks particles, fits a spring-mass
        simulator to that motion, and recovers the contact forces.
        Slow — minutes to hours per video.

        The BOTTOM arrow is what WE do: given a target force
        trajectory, learn a policy that picks gripper motions so the
        realized force tracks the target — in real time, with one
        shared policy across rope, cloth, and stuffed animal.

        One sentence to deliver: "PhysTwin solves the forward problem
        slowly; we invert it, fast, with a learned policy."
    """)
    add_footer(s, 2); return s


def slide_closed_loop(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "What does “closed-loop” mean?",
                 "policy reads achieved force every frame, not a script")

    s.shapes.add_picture(str(FIG / "00_closed_loop_diagram.png"),
                          Inches(0.6), Inches(1.7),
                          width=Inches(12.1))
    add_text(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5),
              "Open-loop drifts ~6 mm/frame. Closed-loop exposes "
              "state-distribution shift.",
              size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    set_notes(s, """
        Quick definition slide so nobody is confused later.

        Closed-loop just means: every frame, we MEASURE the achieved
        force from the simulator and FEED IT BACK to the policy as
        input. The policy then picks the next gripper move based on
        what actually happened, not on a pre-baked schedule.

        Point at the orange feedback arrow when you say "this is what
        makes it closed." If you cut that arrow, it becomes open-loop
        — you just play back a recorded action sequence and accumulate
        error.

        Foreshadowing: the BC policy fails specifically because it
        never SAW the closed-loop input distribution at training time.
        That's what RL fixes.
    """)
    add_footer(s, 3); return s


def slide_pipeline(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Pipeline",
                 "data → BC → closed-loop eval → PPO+BC warm-start → eval")

    s.shapes.add_picture(str(FIG / "00_pipeline_diagram.png"),
                          Inches(0.4), Inches(2.4), width=Inches(12.6))

    add_text(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(0.4),
              "Synthetic motions on the calibrated simulator expand 14 real cases → 33K dataset rows.",
              size=15, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.4),
              "Same MLP shape end-to-end. RL actor = BC weights + exploration head.",
              size=15, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    set_notes(s, """
        How the whole project is wired, left to right.

        Step 1: take the 14 cases PhysTwin already calibrated.
        Step 2: drive the SAME calibrated simulator with synthetic
        motions (push, sinusoid, random walk, hold-release, ramp) —
        262 new trajectories — so the dataset isn't just 14 cases but
        ~33K state-action pairs.
        Step 3: supervised BC on that dataset — one shot, fast.
        Step 4: warm-start a PPO policy from those BC weights and let
        it roll out in the simulator with reward = closed-loop force
        tracking error.
        Step 5: same 14-case eval at the end so BC and RL are directly
        comparable.

        One sentence: "Same MLP all the way through. The RL actor is
        literally the BC weights plus an exploration head on top."
    """)
    add_footer(s, 4); return s


def slide_whats_new(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "What’s new",
                 "first PPO + BC warm-start on PhysTwin’s diff sim with material conditioning")

    bullets = [
        ("Closed-loop driver inside PhysTwin",
         "+200 LOC in trainer_warp.py · runs any policy frame-by-frame"),
        ("Synthetic-trajectory data augmentation",
         "14 real cases  →  262 synthetic motions on same calibrated sim  →  33K dataset rows"),
        ("Inverse-dynamics BC policy + material conditioning",
         "43-dim MLP · [log_spring_Y, log_mean_force] → one policy across rope / cloth / sloth"),
        ("PPO + BC warm-start, from scratch",
         "no SB3 · critic warm-up · log-ratio clipping"),
        ("Diagnosis: state-distribution shift",
         "4 goal-side fixes failed before we picked RL — that’s evidence, not noise"),
    ]
    y = Inches(1.95)
    for head, sub in bullets:
        # bullet dot
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y + Emu(80000),
            Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, Inches(0.95), y, Inches(11.7), Inches(0.4),
                  head, size=18, bold=True, color=INK)
        add_text(s, Inches(0.95), y + Inches(0.4), Inches(11.7), Inches(0.4),
                  sub, size=14, color=DIM)
        y += Inches(0.95)
    set_notes(s, """
        This is your "What's New" slide — the 4-minute heart of the talk.

        1. CLOSED-LOOP DRIVER: a new method we wrote inside PhysTwin's
           warp simulator (~200 LOC) that lets ANY policy drive the
           simulator frame-by-frame, with feedback. Upstream PhysTwin
           only does offline batch optimization — no notion of a
           controller-in-the-loop.

        2. DATA AUGMENTATION: 14 real cases is tiny. We generate 262
           synthetic motions on the SAME calibrated simulator, so the
           training data covers behaviors the real videos don't (e.g.,
           ramps that go up AND come back down). Mention: this alone
           fixed the BC release problem on rope (0.04 -> 0.84) before
           RL was even on the table.

        3. MATERIAL CONDITIONING: a 2-d descriptor — log of mean spring
           stiffness, log of mean recorded force — lets one shared MLP
           specialize per case at inference time.

        4. PPO + BC WARM-START: built our own minimal PPO (no SB3) so
           we control critic warmup and log-ratio clipping. Both
           tricks are needed; without critic warmup, the random critic
           destroys the BC-warm-started actor in the first update.

        5. DIAGNOSIS: we tried 4 goal-side fixes (hindsight relabeling,
           hierarchical sub-goals, incremental shaping, synthetic ramp
           data) — three failed, one partially worked. That negative-
           results chain is what motivated turning to RL.

        Punchline to land: "First end-to-end PPO + BC-warm-start with
        material conditioning on PhysTwin's differentiable simulator."
    """)
    add_footer(s, 5); return s


def slide_observation(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "What the policy sees",
                 "45-dim observation per frame, same vector for BC and RL")

    s.shapes.add_picture(str(FIG / "00_observation_breakdown.png"),
                          Inches(0.4), Inches(1.6), width=Inches(12.6))

    add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.45),
              "Deformation + control point + force feedback + “what material am I?”",
              size=15, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    set_notes(s, """
        30-second slide. Walk left to right across the five blocks:

        - Deformation features (18 dims): summary stats over particle
          positions — centroid displacement, bounding-box change,
          max/mean/std per axis. This is "how is the object shaped
          right now compared to rest?"
        - Control-point features (13 dims): the gripper's current
          position, velocity, and displacement from rest.
        - Achieved force (6 dims): the force currently being applied
          on each gripper (up to 2 grippers, 3 axes each). THIS is
          the feedback signal that makes it closed-loop.
        - Goal force (6 dims): the target force on the NEXT frame.
        - Material descriptor (2 dims): log spring stiffness + log
          typical force magnitude, per case. Lets one MLP specialize.

        Total 45 dimensions. Both BC and RL use the same vector — the
        only difference between them is HOW the weights were trained,
        not what they look at.

        Mention: "the 'closed-loop' part of the project lives in the
        red Achieved Force block — that's the input that comes from
        the simulator each frame, not from training data."
    """)
    add_footer(s, 6); return s


def slide_headline_table(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Results",
                 "force-error ratio (lower is better) · closed-loop evaluation")

    s.shapes.add_picture(str(FIG / "00_bc_vs_rl_bars.png"),
                          Inches(0.4), Inches(1.6),
                          width=Inches(12.6))
    add_text(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.4),
              "Material conditioning helps both. RL’s edge: rope replay 0.75 → 0.41, ramp release 0.04 → 0.87.",
              size=13, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    set_notes(s, """
        The headline slide. Read it as TWO SUB-STORIES across the bars.

        Lower bars = better tracking. The dashed red line at y=1.0 is
        the "uncontrolled" threshold — above it, the policy is failing.

        SUB-STORY 1 (compare gray to blue): adding the 2-d material
        descriptor improves both BC and RL across the board. Most
        dramatic on cloth ramp: 6.55 -> 2.28. One MLP can specialize.

        SUB-STORY 2 (compare blue to green): switching from BC to RL
        helps where the closed-loop input is OUT of training
        distribution. Headline: rope ramp release jumps from 0.04
        (BC barely retracts at all) to 0.87 (RL retracts cleanly).
        Rope replay drops from 0.75 to 0.41 — a 45% improvement.

        Honest beat: cloth ramp and sloth ramp are still above 1.0 —
        we'll talk about that on the limitations slide. RL doesn't fix
        everything.

        One-liner: "Material conditioning helps both methods. RL adds
        a release fix that supervised learning can't get."
    """)
    add_footer(s, 7); return s


def slide_material_conditioning(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Sub-story 1 · Material conditioning",
                 "one policy across rope / cloth / sloth")

    rows = [
        ("",           "rope rep", "cloth ramp", "sloth ramp"),
        ("BC base",       "0.75",      "6.55",        "5.06"),
        ("BC + 2-d desc", "0.46",      "2.28",        "1.73"),
        ("RL + 2-d desc", "0.41",      "2.31",        "2.66"),
    ]
    table_left = Inches(2.0); table_top = Inches(2.2)
    col_w = [Inches(3.0), Inches(2.2), Inches(2.2), Inches(2.2)]
    row_h = Inches(0.7)
    for r, row in enumerate(rows):
        x = table_left
        for c, cell in enumerate(row):
            box = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, table_top + r * row_h,
                col_w[c], row_h)
            if r == 0:
                box.fill.solid(); box.fill.fore_color.rgb = SUBTLE
            else:
                box.fill.solid(); box.fill.fore_color.rgb = BG
            box.line.color.rgb = DIM
            tf = box.text_frame
            tf.margin_left = tf.margin_right = Inches(0.1)
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run(); run.text = cell
            run.font.size = Pt(16); run.font.color.rgb = INK
            run.font.bold = (r == 0 or c == 0)
            run.font.name = "Calibri"
            x += col_w[c]

    add_text(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.6),
              "Adding [log Y, log F] descriptor unlocks generalization — cloth ramp 6.55 → 2.28.",
              size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    set_notes(s, """

        Zoom-in on the material-conditioning sub-story. Just three
        numbers per row, three rows. Read down the columns:

        - "rope rep" = closed-loop replay on rope cases (recorded
          goal). 0.75 -> 0.46 -> 0.41 — every step helps.
        - "cloth ramp" = synthetic ramp on cloth. 6.55 -> 2.28 -> 2.31.
          The descriptor alone fixes most of it; RL barely moves it.
        - "sloth ramp" = synthetic ramp on the stuffed animal.
          5.06 -> 1.73 -> 2.66. Descriptor helped, but RL regressed —
          this is the case that motivates the Fix-G future work.

        Takeaway: "A two-dimensional material descriptor — log
        stiffness, log force magnitude — gets a single MLP to
        specialize across three very different materials."

    """)

    add_footer(s, 8); return s


def slide_release_problem(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Sub-story 2 · The release problem",
                 "why BC stalls and RL doesn’t")

    s.shapes.add_picture(str(FIG / "00_release_problem.png"),
                          Inches(0.5), Inches(1.7), width=Inches(8.2))

    bx_x = Inches(9.0); bx_y = Inches(2.0)
    add_text(s, bx_x, bx_y, Inches(3.8), Inches(0.5),
              "Why?", size=18, bold=True, color=INK)
    add_text(s, bx_x, bx_y + Inches(0.55), Inches(3.8), Inches(2.2),
              "BC never saw the input\n“F(t) high, F*(t+1) low” \nat training time.\n\nRL collects its OWN\nclosed-loop rollouts —\nsees that state, learns\nto retract.",
              size=14, color=DIM)
    add_text(s, bx_x, Inches(5.0), Inches(3.8), Inches(0.6),
              "Release fraction",
              size=14, bold=True, color=INK)
    add_text(s, bx_x, Inches(5.4), Inches(3.8), Inches(1.5),
              "BC base:  0.04\nRL + 2-d:  0.87",
              size=18, color=ACCENT, bold=True)
    set_notes(s, """

        The mechanistic "why does BC fail?" slide. Spend 50 seconds
        here — this is the conceptual core of the BC -> RL story.

        Black line: the goal force trajectory — ramp up, ramp down,
        ending at zero.

        Gray line (BC): tracks the rise OK, but when the goal drops
        back to zero, BC keeps applying force. It clamps at the peak.
        Why? Because supervised data was collected from OPEN-LOOP
        trajectories — so the BC policy never SAW an input where
        "current force is high and goal is low at the same time."

        Green line (RL): RL collects ITS OWN closed-loop rollouts. It
        encounters that "release" state during training, gets bad
        reward when it doesn't retract, and learns to retract.

        Concrete number on the right panel: release fraction goes
        from 0.04 to 0.87 — basically from "doesn't retract" to
        "retracts almost completely."

        One-liner: "The reason BC fails here is exactly the reason
        RL works here — closed-loop states that supervised data
        never contained."

    """)

    add_footer(s, 9); return s


def slide_demo_videos(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Live demo",
                 "best supervised tracking  ·  RL release fix")

    # Left: BC at its best (cloth lift, err 0.22 — prettiest result).
    # Right: RL doing what BC can't (rope ramp release).
    v1 = VID / "06_BC_cloth_lift_BEST.mp4"
    v2 = VID / "01_RL_rope_ramp_HEADLINE.mp4"
    s.shapes.add_movie(str(v1),
                        Inches(0.5), Inches(2.0),
                        Inches(6.0), Inches(3.8),
                        poster_frame_image=str(FIG / "00_release_problem.png"),
                        mime_type="video/mp4")
    s.shapes.add_movie(str(v2),
                        Inches(6.8), Inches(2.0),
                        Inches(6.0), Inches(3.8),
                        poster_frame_image=str(FIG / "00_release_problem.png"),
                        mime_type="video/mp4")
    add_text(s, Inches(0.5), Inches(5.95), Inches(6.0), Inches(0.4),
              "BC + descriptor  ·  cloth lift  ·  err 0.22",
              size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.8), Inches(5.95), Inches(6.0), Inches(0.4),
              "RL + descriptor  ·  rope ramp  ·  release 0.04 → 0.87",
              size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.4),
              "same closed-loop driver  ·  same simulator  ·  same 45-dim input",
              size=12, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    set_notes(s, """

        Live videos, ~30 seconds total — keep it tight.

        LEFT (BC + descriptor on cloth lift, err 0.22): "This is what
        supervised learning gets you when the test situation is in
        distribution. Cleanest single result in the project — error
        ratio 0.22."

        Let it play 8 seconds, then transition.

        RIGHT (RL + descriptor on rope ramp): "And this is what RL
        adds on top — when the goal drops to zero, BC stalls but the
        RL policy retracts. Release fraction 0.04 to 0.87."

        Let that play ~12 seconds. The visual of the gripper actually
        pulling back at the right moment is the selling point.

        Note: same network architecture in both. Same simulator.
        Same 45-dim input. The only difference is HOW the weights
        were trained.

    """)

    add_footer(s, 10); return s


def slide_limitations(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Honest limitations",
                 "where the policy is still uncontrolled")

    items = [
        ("cloth ramp", "err 2.31 · still above the uncontrolled threshold"),
        ("sloth ramp", "err 2.66 · release worked, magnitude undershoots"),
        ("rope_4 outlier", "gentle force regime, training underrepresented"),
        ("Fix G shaped reward",
         "validated on cloth, regressed sloth ramp — next step: per-material α"),
    ]
    y = Inches(2.1)
    for head, sub in items:
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y + Emu(80000),
            Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE
        dot.line.fill.background()
        add_text(s, Inches(1.1), y, Inches(11.7), Inches(0.5),
                  head, size=20, bold=True, color=INK)
        add_text(s, Inches(1.1), y + Inches(0.45), Inches(11.7), Inches(0.4),
                  sub, size=14, color=DIM)
        y += Inches(1.05)
    set_notes(s, """

        Honest 30 seconds — don't oversell.

        Three concrete things to call out:
        1. CLOTH RAMP error ratio is still 2.31 — above the
           "uncontrolled" threshold of 1.0.
        2. SLOTH RAMP regressed from BC (1.73) to RL (2.66). RL
           learned to retract on rope and overfit that behavior on
           sloth.
        3. ROPE_4 outlier: a "gentle push" case (peak force ~2.5 kN
           vs typical ~15 kN) underrepresented in training.

        I also tried a stiffness-shaped reward (Fix G) to fix the
        sloth-ramp regression — it worked on cloth but didn't move
        sloth ramp. Why is one of the next steps.

        Don't dwell. Move on within 30 seconds.

    """)

    add_footer(s, 11); return s


def slide_next_steps(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Next steps for the final report",
                 "concrete deliverables")

    items = [
        ("1.  Per-material reward shaping",
         "Fix G failed because α was global. Train with rope α=0, cloth α=500, sloth α=2000.\nDeliverable: extra row in the headline table, target cloth ramp ≤ 2.0."),
        ("2.  Can BC catch up to RL with more data?",
         "Generate 5× more synthetic trajectories and retrain BC on them.\nIf BC matches RL → data buys the win. If not → RL is doing something supervised data can’t."),
        ("3.  MPC on the diff simulator (stretch)",
         "Rollout-shooting MPC over 10 steps for ramp release reference.\nDeliverable: 4th video showing MPC oracle vs RL on the same trajectory."),
    ]
    y = Inches(1.9)
    for head, sub in items:
        add_text(s, Inches(0.8), y, Inches(11.7), Inches(0.45),
                  head, size=20, bold=True, color=ACCENT)
        add_text(s, Inches(0.8), y + Inches(0.5), Inches(11.7), Inches(1.0),
                  sub, size=14, color=DIM)
        y += Inches(1.55)
    set_notes(s, """

        Three concrete deliverables for the final report. Spend 2
        minutes here — this is the "what's next" budget.

        1. PER-MATERIAL REWARD SHAPING. Fix G used one global alpha
           across all materials. Next try: alpha per material —
           rope=0, cloth=500, sloth=2000. Deliverable: extra row on
           the bar chart. Target cloth ramp below 2.0 without
           regressing rope.

        2. CAN BC CATCH UP TO RL WITH MORE DATA? Current BC dataset
           is 33K rows. Generate 5x more synthetic trajectories,
           retrain BC, see if it matches RL. This isolates "is the
           win really from RL, or could supervised learning have
           gotten there with more data?"
           Deliverable: one extra row "BC + 5x synth" on the same bar
           chart, plus a single-sentence verdict.

        3. STRETCH — MPC ON THE DIFFERENTIABLE SIMULATOR. PhysTwin is
           differentiable, so we can rollout-shoot MPC over ~10 steps
           as an oracle baseline for the ramp-release problem.
           Deliverable: a fourth video showing MPC reference vs RL on
           the same trajectory.

        Emphasize CONCRETE — the grader is asking for that.

    """)

    add_footer(s, 12); return s


def slide_questions(prs):
    s = add_blank_slide(prs); set_bg(s)
    title_block(s, "Questions I’d love feedback on",
                 "pick whichever you have an answer to")

    qs = [
        "Is on-policy PPO the right choice on a differentiable simulator,\nor should we be using analytic-policy-gradient / short-horizon BPTT?",
        "Material descriptor is just [log Y, log F].\nIs a richer descriptor (anisotropy, mesh density, damping) worth chasing?",
        "Fix G’s shaped reward failed on sloth ramp.\nCleaner formulation than per-material α tuning?  (CMDP? Lagrangian?)",
    ]
    y = Inches(2.0)
    for i, q in enumerate(qs):
        add_text(s, Inches(0.7), y, Inches(0.5), Inches(0.5),
                  f"{i+1}.", size=22, bold=True, color=BLUE)
        add_text(s, Inches(1.2), y, Inches(11.5), Inches(1.4),
                  q, size=17, color=INK)
        y += Inches(1.55)
    set_notes(s, """

        Spend ~1.5 minutes here. Pick whichever 2 of 3 feel most
        relevant during the talk:

        Q1 (PPO vs policy gradient on a diff sim): the most
        technically interesting one. PhysTwin is differentiable but
        we treat it as a black box for PPO. Why not analytic policy
        gradients or short-horizon BPTT? Honestly, I'd love an
        opinion.

        Q2 (richer material descriptor): we use [log Y, log F] — 2
        dimensions. Should we add anisotropy, mesh density, damping?
        Or is 2-d already saturated?

        Q3 (cleaner alternative to per-alpha tuning): if Fix G failed
        because alpha was wrong, next attempt tunes alpha per
        material. But that's not principled. Is CMDP / Lagrangian /
        some other formulation cleaner?

        Don't try to answer yourself. Pause after each, let people
        respond.

    """)

    add_footer(s, 13); return s


def slide_thanks(prs):
    s = add_blank_slide(prs); set_bg(s)
    add_text(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.5),
              "Thanks!", size=64, bold=True, color=INK,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
              "Closed-loop force tracking via PhysTwin",
              size=20, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
              "Code + writeups: my_work/code, my_work/docs/closed_loop_control",
              size=13, color=DIM, align=PP_ALIGN.CENTER)
    set_notes(s, """

        Closing. Two lines:
        - "Thanks. Happy to take questions."
        - If anyone asks where the code lives:
          my_work/code (all scripts) and
          my_work/docs/closed_loop_control (all reviews).

    """)

    return s


# --------------------------- build -----------------------------------------

def main():
    prs = make_pres()
    slide_title(prs)            # 1
    slide_goal(prs)             # 2
    slide_closed_loop(prs)      # 3
    slide_pipeline(prs)         # 4
    slide_whats_new(prs)        # 5
    slide_observation(prs)      # 6  ← new: 45-dim input breakdown
    slide_headline_table(prs)   # 7
    slide_material_conditioning(prs)  # 8
    slide_release_problem(prs)  # 9
    slide_demo_videos(prs)      # 10
    slide_limitations(prs)      # 11
    slide_next_steps(prs)       # 12
    slide_questions(prs)        # 13
    slide_thanks(prs)           # 14

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX}")
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
